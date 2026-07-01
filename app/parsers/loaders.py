from __future__ import annotations

import csv
import json
import re
import xml.etree.ElementTree as ET
from pathlib import Path

from pypdf import PdfReader

from app.config import settings

# Raster images — parsed for text + visual content via the vision model (OCR).
IMAGE_EXTENSIONS = {
    ".jpg",
    ".jpeg",
    ".png",
    ".gif",
    ".webp",
    ".bmp",
    ".tif",
    ".tiff",
}

SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".pdf",
    ".docx",
    ".xlsx",
    ".csv",
    ".pptx",
    ".xer",
    ".xml",
    ".ifc",
    ".gltf",
    ".glb",
} | IMAGE_EXTENSIONS

BINARY_EXTENSIONS = {
    ".rvt",
    ".dwg",
    ".dxf",
    ".nwd",
    ".nwc",
    ".fbx",
    ".obj",
    ".3dm",
    ".msg",
}

# Accepted for upload; indexed with a placeholder summary until converted.
STORE_ONLY_EXTENSIONS = BINARY_EXTENSIONS


def load_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def load_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    page_texts: list[str] = []
    for i, page in enumerate(reader.pages):
        if i >= settings.max_pdf_pages:
            page_texts.append(
                f"[… PDF truncated at {settings.max_pdf_pages} pages during extraction …]"
            )
            break
        page_texts.append(page.extract_text() or "")
    return "\n\n".join(page_texts).strip()


def load_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                parts.append(line)
    return "\n".join(parts)


def load_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = [f"Spreadsheet: {path.name}"]
    max_rows_per_sheet = 1000
    for ws in wb.worksheets:
        parts.append(f"\n## Sheet: {ws.title}")
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows_per_sheet:
                parts.append(f"... (rows beyond {max_rows_per_sheet} omitted)")
                break
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def load_csv(path: Path) -> str:
    parts: list[str] = [f"CSV file: {path.name}"]
    max_rows = 5000
    with path.open("r", newline="", encoding="utf-8", errors="replace") as fh:
        sample = fh.read(8192)
        fh.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
        except csv.Error:
            dialect = csv.excel
        reader = csv.reader(fh, dialect)
        for i, row in enumerate(reader):
            if i >= max_rows:
                parts.append(f"... (rows beyond {max_rows:,} omitted)")
                break
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def load_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = [f"PowerPoint presentation: {path.name}"]
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"\n## Slide {idx}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    line = " | ".join(c for c in cells if c)
                    if line:
                        parts.append(line)
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                parts.append(f"[Speaker notes] {notes}")
    return "\n".join(parts)


def load_image(path: Path) -> str:
    """OCR + visual analysis of an image so it becomes searchable/reviewable."""
    from app.llm import describe_image

    size_kb = path.stat().st_size / 1024
    header = f"Image file: {path.name} ({size_kb:.0f} KB, {path.suffix.lower()})"
    try:
        analysis = describe_image(path)
    except Exception:
        analysis = ""

    if analysis and analysis.strip():
        return f"{header}\n\n{analysis.strip()}"

    return (
        f"{header}\n"
        "Image stored, but automatic text/visual analysis was unavailable "
        "(no AI key configured or the image could not be read). "
        "Ask about it again once analysis is available, or upload a searchable PDF."
    )


def load_xer(path: Path) -> str:
    """Parse Primavera P6 XER export into readable text."""
    lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
    sections: dict[str, list[str]] = {}
    current_table: str | None = None
    fields: list[str] = []

    for line in lines:
        if not line.startswith("%"):
            continue
        parts = line.split("\t")
        tag = parts[0]

        if tag == "%T" and len(parts) > 1:
            current_table = parts[1].strip()
            sections.setdefault(current_table, [])
            fields = []
        elif tag == "%F" and current_table:
            fields = [p.strip() for p in parts[1:]]
        elif tag == "%R" and current_table and fields:
            values = [p.strip() for p in parts[1:]]
            row = dict(zip(fields, values))
            # Keep high-value schedule tables compact but searchable.
            if current_table in {
                "PROJECT",
                "PROJWBS",
                "TASK",
                "TASKPRED",
                "CALENDAR",
                "RSRC",
                "ACTVTYPE",
            }:
                summary = ", ".join(f"{k}={v}" for k, v in row.items() if v)
                sections[current_table].append(summary)

    chunks: list[str] = [f"Primavera P6 schedule export ({path.name})"]
    for table, rows in sections.items():
        if not rows:
            continue
        chunks.append(f"\n## {table} ({len(rows)} rows)")
        chunks.extend(rows[:200])
        if len(rows) > 200:
            chunks.append(f"... ({len(rows) - 200} more rows omitted)")
    return "\n".join(chunks)


def load_xml_schedule(path: Path) -> str:
    """Best-effort text from P6 or generic schedule XML."""
    try:
        root = ET.parse(path).getroot()
    except ET.ParseError:
        return load_text_file(path)[:50000]

    chunks: list[str] = [f"XML schedule/document ({path.name})"]
    for elem in root.iter():
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        if tag.lower() in {
            "activity",
            "task",
            "wbs",
            "project",
            "milestone",
            "resource",
        }:
            attrs = " ".join(f'{k}="{v}"' for k, v in elem.attrib.items())
            text = (elem.text or "").strip()
            line = f"{tag}: {attrs} {text}".strip()
            if line:
                chunks.append(line)
    return "\n".join(chunks[:500])


def load_ifc(path: Path) -> str:
    """Extract names and types from IFC (STEP or XML) without heavy CAD libs."""
    raw = path.read_text(encoding="utf-8", errors="replace")[:2_000_000]
    chunks: list[str] = [f"IFC building model ({path.name})"]

    # STEP physical file (most .ifc)
    for match in re.finditer(
        r"IFC\w+\('([^']*)'[^)]*'([^']*)'[^)]*'([^']*)'",
        raw,
    ):
        entity, gid, name = match.groups()
        if name and name != "":
            chunks.append(f"{entity}: {name} (id={gid})")
        if len(chunks) > 400:
            break

    if len(chunks) <= 2:
        # XML IFC fallback
        for match in re.finditer(r'Name="([^"]+)"', raw):
            chunks.append(f"Element: {match.group(1)}")
            if len(chunks) > 400:
                break

    if len(chunks) <= 2:
        chunks.append(
            "Limited metadata extracted. For richer 3D analysis, export an IFC with named elements or companion PDFs."
        )
    return "\n".join(chunks)


def load_gltf(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".glb":
        return (
            f"glTF binary model ({path.name}, {path.stat().st_size} bytes). "
            "Geometry is not text-indexed; upload PDF/IFC/specs for Q&A, or export a glTF JSON (.gltf) with asset metadata."
        )

    data = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    asset = data.get("asset", {})
    scenes = len(data.get("scenes", []))
    nodes = len(data.get("nodes", []))
    meshes = len(data.get("meshes", []))
    materials = len(data.get("materials", []))
    return (
        f"glTF model ({path.name})\n"
        f"generator: {asset.get('generator', 'unknown')}\n"
        f"version: {asset.get('version', 'unknown')}\n"
        f"scenes: {scenes}, nodes: {nodes}, meshes: {meshes}, materials: {materials}"
    )


def load_binary_placeholder(path: Path) -> str:
    size_mb = path.stat().st_size / (1024 * 1024)
    ext = path.suffix.lower()
    hints = {
        ".rvt": "Export IFC or PDF sheets from Revit for full analysis.",
        ".dwg": "Export PDF or DXF with readable layers for analysis.",
        ".dxf": "Text extraction limited; prefer PDF export for drawings.",
        ".msg": "Outlook message stored; save the email as PDF for full-text search.",
        ".jpg": "Image stored; no OCR yet — upload a searchable PDF for text Q&A.",
        ".jpeg": "Image stored; no OCR yet — upload a searchable PDF for text Q&A.",
        ".png": "Image stored; no OCR yet — upload a searchable PDF for text Q&A.",
    }
    hint = hints.get(ext, "Export PDF, IFC, or schedule XER for best results.")
    return (
        f"Binary construction file stored: {path.name} ({size_mb:.1f} MB, {ext})\n"
        f"{hint}\n"
        "File is on record for comparison questions once companion exports are uploaded."
    )


def load_document(path: Path) -> str | None:
    suffix = path.suffix.lower()

    if suffix in {".txt", ".md"}:
        return load_text_file(path)
    if suffix == ".pdf":
        return load_pdf(path)
    if suffix == ".docx":
        return load_docx(path)
    if suffix == ".xlsx":
        return load_xlsx(path)
    if suffix == ".csv":
        return load_csv(path)
    if suffix == ".pptx":
        return load_pptx(path)
    if suffix in IMAGE_EXTENSIONS:
        return load_image(path)
    if suffix == ".xer":
        return load_xer(path)
    if suffix == ".xml":
        return load_xml_schedule(path)
    if suffix == ".ifc":
        return load_ifc(path)
    if suffix == ".gltf":
        return load_gltf(path)
    if suffix == ".glb":
        return load_gltf(path)
    if suffix in STORE_ONLY_EXTENSIONS:
        return load_binary_placeholder(path)

    return None
